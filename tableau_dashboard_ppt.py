from tableauserverclient import Server
import tableauserverclient as TSC
import os
from pptx import Presentation #
from pptx.util import Inches
import os
from PIL import Image
import glob
import numpy as np
import json
import shutil




def read_config():
    try:
        with open("Tableau_Token_Config.txt") as f:
            data = f.read()
        js = json.loads(data)
        return js
    except:
        print("Config File is not the right format, please fix the format of the config file")
        return False


def authenticate():
    try: 
        server = Server(TABLEAU_SERVER_URL,use_server_version=True)
        server.auth.sign_in(tableau_token)
        return server
    except:
        print("Connection error make sure your connected to the network")
        return False

def filter_dashboards(server,WORKBOOKS_TO_BE_EXPORTED):
    try:
        wb_book = {}
        #workbook_ids = []
        workbook_ids = {}
        for wb_name_new in TSC.Pager(server.workbooks):
            wb_book[wb_name_new.name] = wb_name_new.id

        for k,v in WORKBOOKS_TO_BE_EXPORTED.items():
            if k in wb_book.keys():
                #workbook_ids.append(wb_book[i])
                workbook_ids[wb_book[k]] = v
        print(workbook_ids)
        for k,v in workbook_ids.items():
            print(k,v)
        #print(workbook_ids.values())
        return workbook_ids
    except:
        print("Failed to Filter the required dashboards from the server")
        return False

def export_to_image(server, workbook_id,view_names, filename, filters):
    try:
        imagelist = []
        workbook = server.workbooks.get_by_id(workbook_id)
        print(workbook)
        views = workbook.views
        image_req_option = TSC.ImageRequestOptions(imageresolution=TSC.ImageRequestOptions.Resolution.High, maxage=1)
        if type(filters) is dict:
            for filter,filter_value in filters.items():
                image_req_option.vf(filter, filter_value)
        if len(views) > 0:
            for i in range(len(views)):
                if views[i].name in view_names:
                    view = views[i]  
                    server.views.populate_image(view, image_req_option)
                    with open(filename+views[i].id+".png", "wb") as f:
                        f.write(view.image)
                    print(f"Exported workbook '{filename}' as an image.")
                    imagelist.append(filename+views[i].id+".png")
            return True,imagelist
        else:
            print(f"No views found in workbook '{filename}'.")
            return False
    except:
        print("Failed to Export dashboard as images")
        return False


def exporttoimagesprocess(workbook_ids,slideslist,filters):
    try:
        for workbook_id,view_names in workbook_ids.items():
            filename = os.path.join(EXPORT_DIR, f"{workbook_id}_")
            imgstatus,imagelist = export_to_image(server, workbook_id,view_names, filename,filters)
            slideslist.append(imagelist)
        return True
    except:
        print("Failed to complete the export to images process")
        return False


def writetopowerpoint(orderedimagearr):
    try:
        imgarr = [] #later to be array of images
        image_dir = f".\exported_images"

        #Makes the array of images into imgarr
        for filename in os.listdir(image_dir):
            if filename.endswith(".jpg") or filename.endswith(".png"):
                imgarr.append(os.path.join(image_dir, filename))

        imgarr.reverse()
        #Makes the slides in powerpoint
        prs = Presentation('template\weekly_template.pptx')
        blank_slide_layout = prs.slide_layouts[6]

        for numslides in range(len(orderedimagearr)):
            slide = prs.slides.add_slide(blank_slide_layout)
            left = Inches(0.02)
            top = Inches(0.02)
            #right = Inches(0.5)
            pic = slide.shapes.add_picture(orderedimagearr[numslides], left, top, Inches(13.3),Inches(6.96))

        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[0])
        prs.save('Exported_Dashboards.pptx')
        return True
    except:
        print("Failed to export to PowerPoint")
        return False


    
EXPORT_DIR = "./exported_images"
if os.path.exists(EXPORT_DIR) and os.path.isdir(EXPORT_DIR):
    shutil.rmtree(EXPORT_DIR)    
os.makedirs(EXPORT_DIR, exist_ok=True)

existing_ppt = './Exported_Dashboards.pptx'
if os.path.exists(existing_ppt):
    os.remove(existing_ppt)

config_val = read_config()
print(type(config_val))
print(config_val)
slideslist = []
if config_val:
    print("Config File parsed successfully")

    for site,site_values in config_val.items():
        #print(site_values)
        TABLEAU_SERVER_URL = site_values['tableau_server']
        TOKEN_NAME = site_values['token_name']
        TOKEN_VALUE = site_values['token_value']
        SITE_NAME = site_values['site_name']
        WORKBOOKS_TO_BE_EXPORTED = site_values['workbooks_viewnames']
        if "Filters" in site_values:
            FILTERS = site_values["Filters"]
        else:
            FILTERS = ""
        print(FILTERS)

        tableau_token = TSC.PersonalAccessTokenAuth(TOKEN_NAME,TOKEN_VALUE,SITE_NAME)

        server = authenticate()
        print(server)

        if server:
            print("Successfully Authenticated to the server")
        else:
            print("Failed to authenticate to the server")

        workbook_ids = filter_dashboards(server,WORKBOOKS_TO_BE_EXPORTED)

        if workbook_ids:
            print("Filtered the required dashboards successfully")
        else:
            print("Failed to find the requested dashboards from the server")

        status = exporttoimagesprocess(workbook_ids,slideslist,FILTERS)
        if status:
            print("Export to image process completed")
        else:
            print("Export to image process failed")
    flat_list = [item for sublist in slideslist for item in sublist]
    status1 = writetopowerpoint(flat_list)
    if status1:
        print("Writing the images to PowerPoint done successfully")
    else:
        print("Failed to write the images to PowerPoint")
else:
    print("Failed to parse the config file, Check the format of your config file")





